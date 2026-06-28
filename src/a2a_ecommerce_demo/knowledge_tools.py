from __future__ import annotations

import csv
import json
import json as json_lib
import os
import re
import xml.etree.ElementTree as ET
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any

from src.a2a_ecommerce_demo.xlsx_fallback_tools import read_xlsx_sheets

PROJECT_ROOT = Path(__file__).resolve().parents[2]
WIKI_DIR = Path(os.getenv("A2A_WIKI_DIR", PROJECT_ROOT / "wiki")).resolve()
RAW_DIR = Path(os.getenv("A2A_RAW_DIR", PROJECT_ROOT / "raw")).resolve()
LARGE_EXCEL_BYTES = int(os.getenv("A2A_INTERACTIVE_EXCEL_BYTES", str(10 * 1024 * 1024)))
HUGE_EXCEL_BYTES = 100 * 1024 * 1024

SUPPORTED_RAW_SUFFIXES = {
    ".md",
    ".markdown",
    ".txt",
    ".csv",
    ".tsv",
    ".xlsx",
    ".xlsm",
    ".docx",
    ".pdf",
    ".pptx",
    ".json",
    ".jsonl",
    ".xmind",
    ".html",
    ".htm",
    ".xml",
    ".yaml",
    ".yml",
    ".log",
}

STANDARD_WIKI_FOLDERS = [
    "products",
    "suppliers",
    "inventory",
    "datasets",
    "platform-rules",
    "ad-strategy",
    "data-dictionary",
    "cleaning-rules",
    "decisions",
    "logs",
    "sop",
]


def _ensure_wiki_dir() -> None:
    WIKI_DIR.mkdir(parents=True, exist_ok=True)
    RAW_DIR.mkdir(parents=True, exist_ok=True)
    for folder in STANDARD_WIKI_FOLDERS:
        (WIKI_DIR / folder).mkdir(parents=True, exist_ok=True)


def _safe_wiki_path(path: Path) -> Path:
    resolved = path.resolve()
    allowed_roots = [WIKI_DIR, RAW_DIR]
    if not any(root in [resolved, *resolved.parents] for root in allowed_roots):
        raise ValueError(f"Refusing to access outside wiki/raw directories: {resolved}")
    return resolved


def _resolve_wiki_page_path(path: str) -> Path:
    target = Path(path)
    if target.is_absolute():
        return _safe_wiki_path(target)
    parts = target.parts
    if parts and parts[0] == WIKI_DIR.name:
        target = Path(*parts[1:]) if len(parts) > 1 else Path()
    return _safe_wiki_path(WIKI_DIR / target)


def _markdown_files() -> list[Path]:
    _ensure_wiki_dir()
    return sorted(
        path
        for path in WIKI_DIR.rglob("*.md")
        if ".obsidian" not in path.parts
    )


def _relative(path: Path) -> str:
    try:
        return str(path.relative_to(WIKI_DIR)).replace("\\", "/")
    except ValueError:
        return str(path.relative_to(RAW_DIR)).replace("\\", "/")


def _read_text(path: Path) -> str:
    return _safe_wiki_path(path).read_text(encoding="utf-8")


def _slugify(value: str, fallback: str = "note") -> str:
    slug = re.sub(r"[^a-zA-Z0-9\u4e00-\u9fff_-]+", "-", value).strip("-")
    return slug[:80] or fallback


def _cell_to_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, float):
        return str(int(value)) if value.is_integer() else str(value)
    return str(value).strip()


def _row_values(row: tuple[Any, ...]) -> list[str]:
    return [_cell_to_text(value) for value in row]


def _non_empty_count(values: list[str]) -> int:
    return sum(1 for value in values if value)


def _trim_table(rows: list[list[str]]) -> list[list[str]]:
    if not rows:
        return []
    used_columns = sorted(
        {
            column_index
            for row in rows
            for column_index, value in enumerate(row)
            if value
        }
    )
    if not used_columns:
        return []
    first_column = used_columns[0]
    last_column = used_columns[-1]
    return [row[first_column : last_column + 1] for row in rows]


def _render_table(rows: list[list[str]], max_rows: int = 120) -> str:
    rendered = []
    for row in rows[:max_rows]:
        rendered.append(" | ".join(value.replace("\n", " ") for value in row))
    if len(rows) > max_rows:
        rendered.append(f"... {len(rows) - max_rows} more non-empty rows omitted ...")
    return "\n".join(rendered)


def _extract_text_from_html_fragment(value: str) -> str:
    text = re.sub(r"<[^>]+>", "", value)
    return text.replace("&nbsp;", " ").strip()


def _topic_note_from_json(topic: dict[str, Any]) -> str:
    notes = topic.get("notes", {})
    if isinstance(notes, dict):
        plain = notes.get("plain", {})
        if isinstance(plain, dict):
            content = plain.get("content", "")
            if content:
                return str(content).strip()
    return ""


def _render_xmind_topic_json(topic: dict[str, Any], depth: int = 0) -> list[str]:
    title = str(topic.get("title", "")).strip() or "(untitled)"
    indent = "  " * depth
    labels = topic.get("labels") or []
    label_text = f"  {{labels: {', '.join(str(item) for item in labels)}}}" if labels else ""
    link = str(topic.get("href", "")).strip()
    link_text = f"  {{link: {link}}}" if link else ""
    marker_ids = [str(item.get("markerId", "")).strip() for item in topic.get("markers", []) if str(item.get("markerId", "")).strip()]
    marker_text = f"  {{markers: {', '.join(marker_ids)}}}" if marker_ids else ""
    lines = [f"{indent}- {title}{label_text}{link_text}{marker_text}".rstrip()]
    note = _topic_note_from_json(topic)
    if note:
        for line in note.splitlines():
            if line.strip():
                lines.append(f"{indent}  > {line.strip()}")
    children = topic.get("children", {})
    attached = []
    if isinstance(children, dict):
        attached = children.get("attached") or children.get("topics") or []
    for child in attached:
        if isinstance(child, dict):
            lines.extend(_render_xmind_topic_json(child, depth + 1))
    return lines


def _read_xmind_zen(content: str) -> str:
    data = json.loads(content)
    sheets = data if isinstance(data, list) else [data]
    sections: list[str] = []
    for sheet in sheets:
        if not isinstance(sheet, dict):
            continue
        title = str(sheet.get("title", "")).strip() or "Sheet"
        root = sheet.get("rootTopic", {})
        if not isinstance(root, dict):
            continue
        sections.append(f"# Sheet: {title}")
        sections.append("")
        sections.append(f"## {str(root.get('title', '')).strip() or '(untitled)'}")
        sections.append("")
        for child in (root.get("children", {}).get("attached") or []):
            if isinstance(child, dict):
                sections.extend(_render_xmind_topic_json(child, 0))
        sections.append("")
        sections.append("---")
        sections.append("")
    return "\n".join(sections).strip()


def _xml_text(element: ET.Element | None) -> str:
    return "" if element is None else "".join(element.itertext()).strip()


def _render_xmind_topic_xml(topic: ET.Element, depth: int = 0) -> list[str]:
    title = _xml_text(topic.find("title")) or "(untitled)"
    indent = "  " * depth
    labels = [label.text.strip() for label in topic.findall("./labels/label") if label.text and label.text.strip()]
    label_text = f"  {{labels: {', '.join(labels)}}}" if labels else ""
    href = topic.get("{http://www.w3.org/1999/xlink}href", "") or topic.get("href", "")
    link_text = f"  {{link: {href}}}" if href else ""
    marker_ids = [marker.get("marker-id", "").strip() for marker in topic.findall("./marker-refs/marker-ref") if marker.get("marker-id", "").strip()]
    marker_text = f"  {{markers: {', '.join(marker_ids)}}}" if marker_ids else ""
    lines = [f"{indent}- {title}{label_text}{link_text}{marker_text}".rstrip()]
    note_plain = _xml_text(topic.find("./notes/plain"))
    if not note_plain:
        note_plain = _extract_text_from_html_fragment(_xml_text(topic.find("./notes/html")))
    if note_plain:
        for line in note_plain.splitlines():
            if line.strip():
                lines.append(f"{indent}  > {line.strip()}")
    child_topics = topic.findall("./children/topics/topic")
    for child in child_topics:
        lines.extend(_render_xmind_topic_xml(child, depth + 1))
    return lines


def _read_xmind_file(path: Path, max_chars: int) -> str:
    with zipfile.ZipFile(path) as archive:
        names = set(archive.namelist())
        if "content.json" in names:
            content = archive.read("content.json").decode("utf-8", errors="ignore")
            return _read_xmind_zen(content)[:max_chars]
        if "content.xml" in names:
            xml_text = archive.read("content.xml").decode("utf-8", errors="ignore")
            root = ET.fromstring(xml_text)
            sections: list[str] = []
            for sheet in root.findall("./sheet"):
                title = _xml_text(sheet.find("title")) or "Sheet"
                topic = sheet.find("topic")
                if topic is None:
                    continue
                sections.append(f"# Sheet: {title}")
                sections.append("")
                sections.append(f"## {_xml_text(topic.find('title')) or '(untitled)'}")
                sections.append("")
                for child in topic.findall("./children/topics/topic"):
                    sections.extend(_render_xmind_topic_xml(child, 0))
                sections.append("")
                sections.append("---")
                sections.append("")
            return "\n".join(sections).strip()[:max_chars]
    raise ValueError(f"Unsupported or invalid XMind file: {path.name}")


def _read_excel_file(path: Path, max_chars: int) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("Reading Excel files requires openpyxl. Run pip install -r requirements.txt.") from exc

    file_size = path.stat().st_size
    if file_size >= HUGE_EXCEL_BYTES:
        return (
            f"## Large Excel File\n"
            f"- File: {path.name}\n"
            f"- File size: {round(file_size / 1024 / 1024, 2)} MB\n"
            f"- Status: too large for interactive Obsidian ingestion.\n"
            f"- Safe handling: do not parse the workbook inside a frontend chat request. "
            f"Use a dedicated offline extraction/splitting workflow, then ingest the generated summaries or CSV chunks.\n"
        )[:max_chars]

    is_large = file_size >= LARGE_EXCEL_BYTES
    try:
        if is_large:
            workbook = load_workbook(path, read_only=True, data_only=True)
            formula_workbook = None
            max_scan_rows = 300
            max_scan_cols = 80
            max_preview_rows = 80
        else:
            # Some exported workbooks lose data in openpyxl read_only mode, so use normal mode for normal-sized files.
            workbook = load_workbook(path, read_only=False, data_only=True)
            formula_workbook = load_workbook(path, read_only=False, data_only=False)
            max_scan_rows = 5000
            max_scan_cols = 120
            max_preview_rows = 160
    except Exception as exc:
        sections = []
        for sheet in read_xlsx_sheets(path, max_rows=5000, max_cols=120, max_sheets=12):
            rows = _trim_table(sheet["rows"])
            header_row = next((row for row in rows if _non_empty_count(row) >= 2), rows[0] if rows else [])
            section = [
                f"## Sheet: {sheet['sheet']}",
                f"- File size: {round(path.stat().st_size / 1024 / 1024, 2)} MB",
                f"- Dimensions: {sheet.get('dimension') or f'{len(rows)} rows x {sheet.get('max_column', 0)} columns'}",
                f"- Non-empty rows scanned: {len(rows)}",
                f"- Header candidate: {' | '.join(header_row[:40])}",
                f"- OOXML fallback: openpyxl could not parse workbook styles: {exc}",
                "",
                "### Data Preview",
                _render_table(rows, max_rows=160),
            ]
            sections.append("\n".join(section))
        return "\n\n".join(sections)[:max_chars]
    sections = []

    for sheet_index, sheet in enumerate(workbook.worksheets[:12]):
        formula_sheet = formula_workbook.worksheets[sheet_index] if formula_workbook is not None else None
        raw_rows: list[list[str]] = []
        formula_rows: list[list[str]] = []
        max_row = min(sheet.max_row or 0, max_scan_rows)
        max_col = min(sheet.max_column or 0, max_scan_cols)

        if is_large:
            row_iter = sheet.iter_rows(min_row=1, max_row=max_row, max_col=max_col, values_only=True)
            for values_tuple in row_iter:
                values = _row_values(values_tuple)
                if _non_empty_count(values):
                    raw_rows.append(values)
                    formula_rows.append(values)
        else:
            assert formula_sheet is not None
            for row_index in range(1, max_row + 1):
                values = [
                    _cell_to_text(sheet.cell(row=row_index, column=column_index).value)
                    for column_index in range(1, max_col + 1)
                ]
                formula_values = [
                    _cell_to_text(formula_sheet.cell(row=row_index, column=column_index).value)
                    for column_index in range(1, max_col + 1)
                ]
                if _non_empty_count(values) or _non_empty_count(formula_values):
                    raw_rows.append(values if _non_empty_count(values) else formula_values)
                    formula_rows.append(formula_values)

        if not is_large:
            pass
        elif len(raw_rows) >= max_scan_rows:
            raw_rows.append([f"... large workbook preview stopped at {max_scan_rows} rows; use data_cleaning_agent for CSV extraction ..."])

        if not is_large:
            # Filled in by the normal-mode loop above.
            pass

        if not is_large:
            # Keep compatibility with the pre-existing normal workbook formula scan.
            rows_for_formula = formula_rows
        else:
            rows_for_formula = []

        rows = _trim_table(raw_rows)
        formula_rows = _trim_table(rows_for_formula)
        header_row = next((row for row in rows if _non_empty_count(row) >= 2), rows[0] if rows else [])
        formula_cells = []
        for row_index, row in enumerate(formula_rows[:200], start=1):
            for column_index, value in enumerate(row[:80], start=1):
                if value.startswith("="):
                    formula_cells.append(f"R{row_index}C{column_index}: {value}")
                    if len(formula_cells) >= 20:
                        break
            if len(formula_cells) >= 20:
                break

        section = [
            f"## Sheet: {sheet.title}",
            f"- File size: {round(path.stat().st_size / 1024 / 1024, 2)} MB",
            f"- Dimensions: {sheet.max_row} rows x {sheet.max_column} columns",
            f"- Non-empty rows scanned: {len(rows)}",
            f"- Header candidate: {' | '.join(header_row[:40])}",
        ]
        if is_large:
            section.append("- Large-file safe mode: only a preview was extracted for Obsidian to avoid freezing the frontend.")
        if hasattr(sheet, "merged_cells") and sheet.merged_cells.ranges:
            section.append(f"- Merged ranges: {', '.join(str(item) for item in list(sheet.merged_cells.ranges)[:20])}")
        if formula_cells:
            section.append("- Formula cells:")
            section.extend(f"  - {item}" for item in formula_cells)
        section.append("")
        section.append("### Data Preview")
        section.append(_render_table(rows, max_rows=max_preview_rows))
        sections.append("\n".join(section))

    return "\n\n".join(sections)[:max_chars]


def _read_raw_file(path: Path, max_chars: int = 50000) -> str:
    target = _safe_wiki_path(path)
    suffix = target.suffix.lower()
    if suffix in {".md", ".markdown", ".txt", ".yaml", ".yml", ".log"}:
        return target.read_text(encoding="utf-8", errors="ignore")[:max_chars]
    if suffix in {".csv", ".tsv"}:
        delimiter = "\t" if suffix == ".tsv" else ","
        with target.open("r", encoding="utf-8-sig", newline="") as file:
            rows = list(csv.reader(file, delimiter=delimiter))
        preview = "\n".join([", ".join(row) for row in rows[:80]])
        return preview[:max_chars]
    if suffix in {".xlsx", ".xlsm"}:
        return _read_excel_file(target, max_chars)
    if suffix == ".docx":
        try:
            from docx import Document
        except ImportError as exc:
            raise RuntimeError("Reading DOCX files requires python-docx. Run pip install python-docx.") from exc
        document = Document(str(target))
        parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table_index, table in enumerate(document.tables[:20], start=1):
            parts.append(f"## Table {table_index}")
            for row in table.rows[:80]:
                parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        return "\n".join(parts)[:max_chars]
    if suffix == ".pdf":
        try:
            from pypdf import PdfReader
        except ImportError as exc:
            raise RuntimeError("Reading PDF files requires pypdf. Run pip install pypdf.") from exc
        reader = PdfReader(str(target))
        pages = []
        for page_index, page in enumerate(reader.pages[:30], start=1):
            pages.append(f"## Page {page_index}\n{page.extract_text() or ''}")
        return "\n\n".join(pages)[:max_chars]
    if suffix == ".pptx":
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("Reading PPTX files requires python-pptx. Run pip install -r requirements.txt.") from exc
        presentation = Presentation(str(target))
        slides = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            lines = [f"## Slide {slide_index}"]
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    lines.append(text.strip())
                if getattr(shape, "has_table", False):
                    table = getattr(shape, "table", None)
                    if table is not None:
                        for row in table.rows:
                            lines.append(" | ".join(cell.text.strip() for cell in row.cells))
            if slide.has_notes_slide:
                notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
                notes = getattr(notes_frame, "text", "")
                if notes and notes.strip():
                    lines.append(f"Notes: {notes.strip()}")
            slides.append("\n".join(lines))
        return "\n\n".join(slides)[:max_chars]
    if suffix in {".json", ".jsonl"}:
        text = target.read_text(encoding="utf-8", errors="ignore")
        if suffix == ".json":
            try:
                parsed = json_lib.loads(text)
                return json_lib.dumps(parsed, ensure_ascii=False, indent=2)[:max_chars]
            except json_lib.JSONDecodeError:
                return text[:max_chars]
        lines = []
        for index, line in enumerate(text.splitlines()):
            if index >= 100:
                break
            try:
                lines.append(json_lib.dumps(json_lib.loads(line), ensure_ascii=False))
            except json_lib.JSONDecodeError:
                lines.append(line)
        return "\n".join(lines)[:max_chars]
    if suffix == ".xmind":
        return _read_xmind_file(target, max_chars)
    if suffix in {".html", ".htm", ".xml"}:
        try:
            from bs4 import BeautifulSoup
        except ImportError as exc:
            raise RuntimeError("Reading HTML/XML files requires beautifulsoup4. Run pip install -r requirements.txt.") from exc
        text = target.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(text, "html.parser" if suffix in {".html", ".htm"} else "xml")
        for element in soup(["script", "style"]):
            element.decompose()
        return soup.get_text("\n", strip=True)[:max_chars]
    raise ValueError(f"Unsupported raw file type: {target.name}")


def _classify_wiki_folder(filename: str, content: str) -> str:
    haystack = f"{filename}\n{content}".lower()
    rules = [
        ("logs", ["进销存", "台账", "库存表", "库存台账", "出库总量", "入库总量", "期初总量"]),
        ("suppliers", ["supplier", "供应商", "工厂", "交期", "采购", "合同"]),
        ("products", ["sku", "asin", "产品", "商品", "listing", "卖点"]),
        ("sop", ["sop", "流程", "规则", "政策", "补货", "决策标准"]),
        ("platform-rules", ["amazon", "平台", "合规", "违规", "policy", "规则"]),
        ("ad-strategy", ["广告", "acos", "roas", "关键词", "投放", "campaign"]),
        ("decisions", ["复盘", "决策", "会议", "结论", "review"]),
    ]
    for folder, keywords in rules:
        if any(keyword in haystack for keyword in keywords):
            return folder
    return "logs"


def _summarize_raw_content(content: str, max_lines: int = 18) -> str:
    lines = [line.strip() for line in content.splitlines() if line.strip()]
    if not lines:
        return "未提取到可读文本。"
    selected = lines[:max_lines]
    return "\n".join(f"- {line[:220]}" for line in selected)


def _update_index_link(wiki_path: str, title: str) -> None:
    index_path = _safe_wiki_path(WIKI_DIR / "index.md")
    if not index_path.exists():
        index_path.write_text("# A2A Knowledge Wiki\n", encoding="utf-8")
    content = index_path.read_text(encoding="utf-8")
    link = f"- [[{wiki_path.removesuffix('.md')}|{title}]]"
    section = "\n## 自动入库\n"
    if link in content:
        return
    if section not in content:
        content = content.rstrip() + section
    content = content.rstrip() + f"\n{link}\n"
    index_path.write_text(content, encoding="utf-8")


def _score_document(query_terms: list[str], text: str, path: Path) -> int:
    lowered = text.lower()
    path_text = _relative(path).lower()
    score = 0
    for term in query_terms:
        if not term:
            continue
        score += lowered.count(term) * 2
        score += path_text.count(term) * 5
    return score


def _snippet(text: str, terms: list[str], max_chars: int = 500) -> str:
    lowered = text.lower()
    positions = [lowered.find(term) for term in terms if term and lowered.find(term) >= 0]
    if not positions:
        return text[:max_chars].strip()
    start = max(0, min(positions) - 160)
    return text[start : start + max_chars].strip()


def list_wiki_pages() -> str:
    """列出 Obsidian wiki 中的 Markdown 页面。"""
    pages = []
    for path in _markdown_files():
        text = _read_text(path)
        title_match = re.search(r"^#\s+(.+)$", text, re.MULTILINE)
        pages.append(
            {
                "path": _relative(path),
                "title": title_match.group(1).strip() if title_match else path.stem,
                "chars": len(text),
            }
        )
    return json.dumps({"wiki_dir": str(WIKI_DIR), "pages": pages}, ensure_ascii=False, indent=2)


def search_wiki(query: str, limit: int = 5) -> str:
    """按关键词搜索 Obsidian wiki 页面，返回最相关页面和片段。"""
    terms = [term.lower() for term in re.split(r"[\s,，。；;:：/\\|]+", query) if term.strip()]
    results: list[dict[str, Any]] = []
    for path in _markdown_files():
        text = _read_text(path)
        score = _score_document(terms, text, path)
        if score <= 0:
            continue
        results.append(
            {
                "path": _relative(path),
                "score": score,
                "snippet": _snippet(text, terms),
            }
        )
    results.sort(key=lambda item: item["score"], reverse=True)
    return json.dumps({"query": query, "results": results[:limit]}, ensure_ascii=False, indent=2)


def read_wiki_page(path: str) -> str:
    """读取指定 wiki Markdown 页面。path 使用相对路径，例如 sop/restock-policy.md。"""
    target = _resolve_wiki_page_path(path)
    return json.dumps(
        {
            "path": _relative(target),
            "content": _read_text(target),
        },
        ensure_ascii=False,
        indent=2,
    )


def append_decision_note(title: str, content: str) -> str:
    """把有复用价值的决策结论追加为 Obsidian 决策记录。"""
    _ensure_wiki_dir()
    slug = re.sub(r"[^a-zA-Z0-9\u4e00-\u9fff_-]+", "-", title).strip("-")[:70] or "decision"
    path = _safe_wiki_path(WIKI_DIR / "decisions" / f"{datetime.now().strftime('%Y%m%d-%H%M%S')}-{slug}.md")
    path.parent.mkdir(parents=True, exist_ok=True)
    body = f"# {title}\n\n{content.strip()}\n"
    path.write_text(body, encoding="utf-8")
    return json.dumps({"saved_to": str(path), "wiki_path": _relative(path)}, ensure_ascii=False, indent=2)


def append_dataset_insight(dataset_slug: str, title: str, content: str, insight_type: str = "analysis") -> str:
    """把可复用的数据结论写回对应数据集的 wiki 记忆层。"""
    _ensure_wiki_dir()
    safe_dataset_slug = _slugify(dataset_slug, "dataset")
    safe_insight_type = _slugify(insight_type, "analysis")
    insight_dir = _safe_wiki_path(WIKI_DIR / "datasets" / safe_dataset_slug / "insights")
    insight_dir.mkdir(parents=True, exist_ok=True)
    slug = _slugify(title, "insight")
    filename = f"{datetime.now().strftime('%Y%m%d-%H%M%S')}-{safe_insight_type}-{slug}.md"
    path = _safe_wiki_path(insight_dir / filename)
    body = (
        f"# {title}\n\n"
        f"- Dataset: `{safe_dataset_slug}`\n"
        f"- Insight type: `{safe_insight_type}`\n"
        f"- Recorded at: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
        f"{content.strip()}\n"
    )
    path.write_text(body, encoding="utf-8")
    return json.dumps({"saved_to": str(path), "wiki_path": _relative(path)}, ensure_ascii=False, indent=2)


DURABLE_INSIGHT_TYPES = {
    "field_definition": "字段口径",
    "quality_issue": "质量异常",
    "business_rule": "业务规则",
    "analysis_template": "分析模板",
    "retrospective": "复盘结论",
}


def append_durable_insight(
    dataset_slug: str,
    title: str,
    content: str,
    insight_type: str = "business_rule",
    evidence_paths: list[str] | None = None,
) -> str:
    """只把可长期复用的字段口径、质量异常、业务规则、分析模板和复盘结论写回 wiki。"""
    normalized_type = _slugify(insight_type, "business_rule")
    if normalized_type not in DURABLE_INSIGHT_TYPES:
        return json.dumps(
            {
                "status": "rejected",
                "reason": "Only durable reusable knowledge should be written back to wiki.",
                "requested_type": insight_type,
                "allowed_types": sorted(DURABLE_INSIGHT_TYPES),
            },
            ensure_ascii=False,
            indent=2,
        )

    evidence_paths = evidence_paths or []
    _ensure_wiki_dir()
    safe_dataset_slug = _slugify(dataset_slug, "dataset")
    insight_dir = _safe_wiki_path(WIKI_DIR / "datasets" / safe_dataset_slug / "insights")
    insight_dir.mkdir(parents=True, exist_ok=True)
    slug = _slugify(title, "durable-insight")
    path = _safe_wiki_path(insight_dir / f"{datetime.now().strftime('%Y%m%d-%H%M%S')}-{normalized_type}-{slug}.md")
    evidence_block = "\n".join(f"- `{item}`" for item in evidence_paths if str(item).strip()) or "- 待补充"
    body = (
        "---\n"
        "tags:\n"
        "  - durable-insight\n"
        f"  - {normalized_type}\n"
        f"dataset: {safe_dataset_slug}\n"
        f"insight_type: {normalized_type}\n"
        f"recorded_at: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
        "---\n\n"
        f"# {title}\n\n"
        f"- Dataset: `{safe_dataset_slug}`\n"
        f"- Durable type: `{normalized_type}` ({DURABLE_INSIGHT_TYPES[normalized_type]})\n\n"
        "## Insight\n\n"
        f"{content.strip()}\n\n"
        "## Evidence\n\n"
        f"{evidence_block}\n"
    )
    path.write_text(body, encoding="utf-8")
    return json.dumps(
        {
            "status": "success",
            "saved_to": str(path),
            "wiki_path": _relative(path),
            "insight_type": normalized_type,
            "allowed_types": sorted(DURABLE_INSIGHT_TYPES),
        },
        ensure_ascii=False,
        indent=2,
    )


def list_raw_files() -> str:
    """列出 raw 目录中可入库的原始资料文件。"""
    _ensure_wiki_dir()
    files = []
    for path in sorted(RAW_DIR.rglob("*")):
        if not path.is_file():
            continue
        if path.suffix.lower() not in SUPPORTED_RAW_SUFFIXES:
            continue
        files.append(
            {
                "path": _relative(path),
                "name": path.name,
                "size_bytes": path.stat().st_size,
                "type": path.suffix.lower(),
            }
        )
    return json.dumps({"raw_dir": str(RAW_DIR), "files": files}, ensure_ascii=False, indent=2)


def ingest_raw_file(path: str, title: str = "", target_folder: str = "") -> str:
    """把 raw 目录中的单个资料整理成 Obsidian wiki 页面。"""
    _ensure_wiki_dir()
    raw_path = _safe_wiki_path(RAW_DIR / path)
    content = _read_raw_file(raw_path)
    page_title = title.strip() or raw_path.stem
    folder = target_folder.strip() or _classify_wiki_folder(raw_path.name, content)
    folder = folder.strip("/\\")
    wiki_folder = _safe_wiki_path(WIKI_DIR / folder)
    wiki_folder.mkdir(parents=True, exist_ok=True)

    slug = _slugify(page_title, raw_path.stem)
    wiki_path = _safe_wiki_path(wiki_folder / f"{slug}.md")
    relative_raw = _relative(raw_path)
    extracted_at = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    summary = _summarize_raw_content(content)
    body = f"""---
title: {page_title}
source: {relative_raw}
ingested_at: {extracted_at}
folder: {folder}
tags:
  - raw-ingest
  - {folder}
---

# {page_title}

## Source

- Raw file: `{relative_raw}`
- Ingested at: {extracted_at}
- Suggested folder: `{folder}`

## Summary

{summary}

## Key Facts

- 待人工或 Agent 进一步整理。

## Follow-up Questions

- 这份资料是否影响补货、广告、Listing、供应商或合规决策？
- 是否需要拆分成产品页、SOP 页或供应商页？

## Raw Extract

```text
{content[:6000].strip()}
```
"""
    wiki_path.write_text(body, encoding="utf-8")
    wiki_rel = _relative(wiki_path)
    _update_index_link(wiki_rel, page_title)

    log_title = f"Raw ingest: {page_title}"
    log_content = f"- Raw: `{relative_raw}`\n- Wiki: `[[{wiki_rel.removesuffix('.md')}|{page_title}]]`\n- Time: {extracted_at}\n"
    log_slug = _slugify(log_title, f"Raw-ingest-{raw_path.stem}")
    log_path = _safe_wiki_path(WIKI_DIR / "decisions" / f"{log_slug}.md")
    log_path.parent.mkdir(parents=True, exist_ok=True)
    log_path.write_text(f"# {log_title}\n\n{log_content.strip()}\n", encoding="utf-8")

    return json.dumps(
        {
            "raw_file": relative_raw,
            "wiki_path": wiki_rel,
            "ingest_log_path": _relative(log_path),
            "title": page_title,
            "folder": folder,
            "chars_ingested": len(content),
        },
        ensure_ascii=False,
        indent=2,
    )


def ingest_all_raw_files(limit: int = 10) -> str:
    """批量把 raw 目录中的资料整理成 Obsidian wiki 页面。"""
    _ensure_wiki_dir()
    results = []
    raw_files = [
        path
        for path in sorted(RAW_DIR.rglob("*"))
        if path.is_file() and path.suffix.lower() in SUPPORTED_RAW_SUFFIXES
    ]
    for path in raw_files[:limit]:
        try:
            results.append(json.loads(ingest_raw_file(_relative(path))))
        except Exception as exc:
            results.append({"raw_file": _relative(path), "error": str(exc)})
    return json.dumps({"processed": len(results), "results": results}, ensure_ascii=False, indent=2)
